import fitz
import openai
import dotenv
import os
import re
import json
path = "./PDF/Your papers path .pdf"  ##ここにPDFのパスを入れる
doc = fitz.open(path) 
out = open("Text/output.txt", "wb") 
for page in doc:
    text = page.get_text().encode("utf8") 
    out.write(text)
    out.write(bytes((12,)))
with open("Text/output.txt", "r",encoding="utf-8") as f:
    paper_text = f.read()
with open("Prompts/extract_prompts.txt", "r",encoding="utf-8") as f:
    extract_prompts = f.read()
with open("Prompts/system_prompts.txt", "r",encoding="utf-8") as f:
    system_prompts = f.read()
with open("Prompts/user_format.txt", "r",encoding="utf-8") as f:
    user_prompts = f.read()

user_prompts =user_prompts + paper_text

def _create_message(role,prompt):
    return {
        "role": role,
        "content": prompt
    }
system_message = _create_message("system",extract_prompts)
user_message = _create_message("user",user_prompts)
messages = [system_message,user_message]

dotenv.load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")

response = openai.ChatCompletion.create(
    model = "gpt-4-turbo",
    messages = messages,
    temperature = 0,
)
sys_response = response.choices[0].message['content']

print(sys_response)

def extract_and_save_json(text, output_filename):
    pattern = r"```json\n([\s\S]*?)\n```"
    matches = re.findall(pattern, text)

    if matches:
        json_data = matches[0]
        try:
            json_object = json.loads(json_data)
            for item in json_object:
                if 'contents' in item:
                    item['contents'] = item['contents'].replace('\\n', '\n')
            with open(output_filename, 'w', encoding='utf-8') as f:
                json.dump(json_object, f, ensure_ascii=False, indent=4)
        except json.JSONDecodeError as e:
            print("抽出したデータが有効なJSON形式ではありません。エラー:", e)
    else:
        print("有効なJSONデータがテキスト内に見つかりませんでした。")

extract_and_save_json(sys_response, "json_data/temp.json")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

import json
import re

with open("json_data/temp.json", "r", encoding="utf-8") as f:
    data = json.load(f)

prs = Presentation()
prs.slide_width = Inches(16 * 0.75) 
prs.slide_height = Inches(9 * 0.75)

title_slide_layout = prs.slide_layouts[0]
title_slide = prs.slides.add_slide(title_slide_layout)
title = title_slide.shapes.title
title.left = Inches(1)
title.width = Inches(prs.slide_width.inches - 2) 
title.top = Inches(1.5) 
content = data[0]['contents']
title_text = content.split("\n")[0]
title.text = title_text 
subtitle = title_slide.placeholders[1]
subtitle.text = "Hoge Fuga"  ## ここに発表者の名前を入れる 

for slide_data in data:
    slide = prs.slides.add_slide(prs.slide_layouts[1]) 

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(242, 242, 242)
    title_shape = slide.shapes.title
    title_shape.text = slide_data['title']
    title_shape.width = prs.slide_width - Inches(2) 
    title_shape.height = Inches(1.5)
    title_shape.left = Inches(1) 
    title_shape.top = Inches(0.5) 
    content_shape = slide.placeholders[1]
    content_shape.text = slide_data['contents'].replace("・", "")
    content_shape.width = prs.slide_width - Inches(2) 
    content_shape.height = prs.slide_height - Inches(3) 
    content_shape.left = Inches(1)
    content_shape.top = Inches(1.5)
    content_frame = content_shape.text_frame
    content_frame.margin_top = Inches(0.1)
    content_frame.margin_bottom = Inches(0.1)
    content_frame.margin_left = Inches(0.1)
    content_frame.margin_right = Inches(0.1)
    content_frame.vertical_anchor = MSO_ANCHOR.MIDDLE    
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = slide_data['notes']
    date = datetime.now().strftime("%Y/%m/%d")
    footer_text = f"Organization Name {date}"
    footer_left = (prs.slide_width - Inches(3.0)) / 2  
    footer_top = prs.slide_height - Inches(0.8)
    footer_width = Inches(3.0)
    footer_height = Inches(0.5)
    footer_box = slide.shapes.add_textbox(footer_left, footer_top, footer_width, footer_height)
    footer_frame = footer_box.text_frame
    footer_frame.clear()
    p = footer_frame.paragraphs[0]
    p.text = footer_text
    p.font.bold = True
    p.font.size = Pt(10)

for slide_number, slide in enumerate(prs.slides):
    number_left = prs.slide_width - Inches(1.5)
    number_top = prs.slide_height - Inches(0.8)
    number_width = Inches(1.0)
    number_height = Inches(0.5)

    number_box = slide.shapes.add_textbox(number_left, number_top, number_width, number_height)
    number_frame = number_box.text_frame
    number_frame.clear()
    p = number_frame.paragraphs[0]
    p.text = str(slide_number + 1)
    p.font.bold = True
    p.font.size = Pt(10)
if data:
    title = data[0]['contents']
    safe_title = re.sub(r'[<>:"/\\|?*]', '', title)
    safe_title = safe_title.replace(' ', '_')
    filename = f"{safe_title[:50]}.pptx"
else:
    filename = "presentation.pptx"

# ファイルの保存
prs.save(filename)
print(f"プレゼンテーションが {filename} として保存されました。")
